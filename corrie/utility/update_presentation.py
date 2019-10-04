import os

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt
from pptx.enum.chart import XL_LEGEND_POSITION, XL_DATA_LABEL_POSITION


class UpdatePresentation(object):

    def __init__(self, saved_data, collected_results, current_corrie_file_name):
        self.saved_data = saved_data
        self.collected_results = collected_results
        self.current_corrie_file_name = current_corrie_file_name

    def test1(self):
        # create presentation with 1 slide ------
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])

        # define chart data ---------------------
        chart_data = CategoryChartData()
        chart_data.categories = ['East', 'West', 'Midwest']
        chart_data.add_series('Series 1', (19.2, 21.4, 16.7))

        # add chart to slide --------------------
        x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
        slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data
        )

        pptx_file = self.saved_data['powerpointPath']
        prs.save(pptx_file)

    def create_slides(self):
        prs = Presentation()
        pptx_file = self.saved_data['powerpointPath']

        try:
            f = open(pptx_file, "r+")
            f.close()
        except IOError:
            print('file {} is already open and should be shut before file is saved'.format(pptx_file))
            return

        # create assumptions slide

        slide = prs.slides.add_slide(prs.slide_layouts[1])
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        title_shape.text = 'Assumptions'
        tf = body_shape.text_frame
        p = tf.paragraphs[0]
        p.text = 'Building: ' + self.saved_data['building']
        p = tf.add_paragraph()
        p.text = 'Building front faces: ' + self.saved_data['frontFaces']
        p = tf.add_paragraph()
        p.text = 'Baseline code: ' + self.saved_data['baselineCode']
        p = tf.add_paragraph()
        weather_file_and_path = self.saved_data['weatherPath']
        _, weather_file = os.path.split(weather_file_and_path)
        p.text = 'Weather File: ' + weather_file
        occupancy_areas = self.saved_data['occupancyAreas']
        p = tf.add_paragraph()
        p.text = 'Occupancies'
        for area_name, area_value in occupancy_areas.items():
            p = tf.add_paragraph()
            p.text = '{}: {} sqft'.format(area_name, area_value)
            p.level = 1

        # create individual measure slides
        for slide_name in self.collected_results:
            category_list = []
            series_data = []
            slide_results = self.collected_results[slide_name]
            for option_name in slide_results:
                category_list.append(option_name)
                option_results = slide_results[option_name]
                series_data.append(option_results['net_site_energy'])
            # define chart data ---------------------
            if series_data:
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                shapes = slide.shapes
                shapes.title.text = slide_name

                chart_data = CategoryChartData()
                chart_data.categories = category_list
                print('chart_data.categories: ', chart_data.categories)
                print('tuple(series_data): ', tuple(series_data))
                chart_data.add_series('Series 1', tuple(series_data))

                # add chart to slide --------------------
                x, y, cx, cy = Inches(0.5), Inches(2), Inches(9), Inches(5)
                chart = slide.shapes.add_chart(
                    XL_CHART_TYPE.BAR_CLUSTERED, x, y, cx, cy, chart_data
                ).chart
                value_axis = chart.value_axis
                value_axis_tick_labels = value_axis.tick_labels
                value_axis_tick_labels.number_format = '#,##0'
                value_axis_tick_label_font = value_axis_tick_labels.font
                value_axis_tick_label_font.size = Pt(12)
                value_axis_title = value_axis.axis_title
                value_axis_title.text_frame.text = 'Net Site Energy (kBtu)'
                value_axis.minimum_scale = 0
                print('created slide named: ',slide_name)

        # create pie chart slide
        slide_results = self.collected_results['Aspect Ratio']
        option_results = slide_results['width 1.0 : depth 3.0']
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        shapes = slide.shapes
        shapes.title.text = 'End Use Breakdown'

        chart_data = CategoryChartData()

        total = option_results['net_site_energy']

        end_use_result_labels = ["end_use_heating","end_use_cooling","end_use_interior_lighting","end_use_exterior_lighting",
                                 "end_use_interior_equipment","end_use_exterior_equipment","end_use_fans","end_use_pumps",
                                 "end_use_heat_rejection","end_use_humidification","end_use_heat_recovery","end_use_water_systems",
                                 "end_use_refrigeration", "end_use_generators"]
        end_use_result_labels_used = [label for label in end_use_result_labels if option_results[label] > 0 ]
        end_use_names = [eu.replace('end_use_','').replace('_',' ') for eu in end_use_result_labels_used]
        series_data =  [option_results[label]/total for label in end_use_result_labels_used]
        chart_data.categories = end_use_names
        chart_data.add_series('s1', tuple(series_data))

        x, y, cx, cy = Inches(0.5), Inches(2), Inches(9), Inches(5)

        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data
        ).chart

        #chart.has_legend = True
        #chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        #chart.legend.include_in_layout = False

        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        data_labels.show_category_name = True
        data_labels.number_format = '0%'
        data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END

        prs.save(pptx_file)
